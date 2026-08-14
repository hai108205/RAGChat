"""Document loader — loads files of various formats into text."""

from pathlib import Path


class DocumentLoader:
    """Load documents from file paths, returning raw text."""

    SUPPORTED_EXTENSIONS = frozenset({".pdf", ".docx", ".txt", ".md", ".pptx", ".csv", ".xlsx", ".html"})

    def load(self, file_path: str | Path) -> str:
        """Load a document file and return its full text content.

        Args:
            file_path: Path to the document file.

        Returns:
            Full text content of the document.

        Raises:
            ValueError: If the file format is not supported.
            FileNotFoundError: If the file does not exist.
        """
        file_path = Path(file_path)

        ext = file_path.suffix.lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file format: {ext}")

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        if ext == ".pdf":
            return self._load_pdf(file_path)
        elif ext == ".docx":
            return self._load_docx(file_path)
        elif ext == ".pptx":
            return self._load_pptx(file_path)
        elif ext in {".txt", ".md", ".html"} or ext in {".csv"}:
            return self._load_text(file_path)
        elif ext in {".xlsx"}:
            return self._load_xlsx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    def _load_pdf(self, path: Path) -> str:
        from PyPDF2 import PdfReader

        reader = PdfReader(str(path))
        texts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                texts.append(text)
        return "\n\n".join(texts)

    def _load_docx(self, path: Path) -> str:
        from docx import Document

        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def _load_pptx(self, path: Path) -> str:
        from pptx import Presentation

        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n\n".join(texts)

    def _load_text(self, path: Path) -> str:
        return path.read_text(encoding="utf-8", errors="replace")

    def _load_xlsx(self, path: Path) -> str:
        import openpyxl

        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        texts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append(" | ".join(str(cell) if cell is not None else "" for cell in row))
            texts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(texts)
